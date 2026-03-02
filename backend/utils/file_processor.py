import os
import logging
import fitz
import mammoth
from pptx import Presentation

class FileProcessor:
    def __init__(self, upload_folder='uploads/'):
        self.upload_folder = upload_folder
        if not os.path.exists(self.upload_folder):
            os.makedirs(self.upload_folder)

    def extract_text_from_pdf(self, file_path):
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text

    def extract_text_from_docx(self, file_path):
        with open(file_path, "rb") as docx_file:
            result = mammoth.extract_raw_text(docx_file)
            return result.value

    def extract_text_from_pptx(self, file_path):
        """Extract text from a .pptx PowerPoint file.

        Iterates over every slide and shape, pulling text from text-frames
        (titles, body placeholders, free text-boxes) and table cells.
        """
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            text_parts.append(para_text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                text_parts.append(cell_text)
        return "\n".join(text_parts)

    def process_file(self, file):
        file_path = os.path.join(self.upload_folder, file.filename)
        file.save(file_path)
        content = ""

        try:
            if file.filename.endswith('.txt'):
                with open(file_path, 'r') as f:
                    content = f.read()
            elif file.filename.endswith('.pdf'):
                content = self.extract_text_from_pdf(file_path)
            elif file.filename.endswith('.docx'):
                content = self.extract_text_from_docx(file_path)
            elif file.filename.endswith('.pptx'):
                content = self.extract_text_from_pptx(file_path)
            elif file.filename.endswith('.ppt'):
                logging.warning(
                    "Legacy .ppt format is not supported. "
                    "Please convert to .pptx and try again."
                )
        finally:
            if os.path.exists(file_path):
                os.remove(file_path)

        return content
