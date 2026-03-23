"""Tests for PPTX text extraction in FileProcessor.

Imports the real FileProcessor from Generator.main. Heavy ML dependencies
(torch, transformers, sense2vec, etc.) are mocked at the sys.modules level
so they are not actually loaded during test collection.
"""
import os
import sys
from unittest.mock import MagicMock

import pytest
from pptx import Presentation
from pptx.util import Inches

# ── Mock heavy ML dependencies before importing Generator.main ───────────────
# This prevents model loading while still testing the real FileProcessor code.

_mock_torch = MagicMock()
_mock_torch.cuda.is_available.return_value = False
sys.modules.setdefault("torch", _mock_torch)
sys.modules.setdefault("transformers", MagicMock())
sys.modules.setdefault("sense2vec", MagicMock())
sys.modules.setdefault("google.oauth2", MagicMock())
sys.modules.setdefault("google.oauth2.service_account", MagicMock())
sys.modules.setdefault("googleapiclient", MagicMock())
sys.modules.setdefault("googleapiclient.discovery", MagicMock())
sys.modules.setdefault("en_core_web_sm", MagicMock())

from Generator.main import FileProcessor  # noqa: E402


# ── helpers ──────────────────────────────────────────────────────────────────

def _create_pptx(tmp_path, filename, texts):
    """Create a minimal .pptx with one text-box per item in *texts*."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    for t in texts:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = t
    path = os.path.join(tmp_path, filename)
    prs.save(path)
    return path


def _create_pptx_with_table(tmp_path, filename, rows_data):
    """Create a .pptx containing a single table."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cols = max(len(r) for r in rows_data)
    table_shape = slide.shapes.add_table(
        len(rows_data), cols, Inches(1), Inches(1), Inches(6), Inches(2)
    )
    table = table_shape.table
    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            table.cell(ri, ci).text = cell_text
    path = os.path.join(tmp_path, filename)
    prs.save(path)
    return path


# ── tests ────────────────────────────────────────────────────────────────────

class TestExtractTextFromPptx:

    def test_simple_text(self, tmp_path):
        path = _create_pptx(tmp_path, "simple.pptx", ["Hello World", "Second box"])
        fp = FileProcessor(upload_folder=str(tmp_path))
        result = fp.extract_text_from_pptx(path)
        assert "Hello World" in result
        assert "Second box" in result

    def test_table_text(self, tmp_path):
        rows = [["Name", "Score"], ["Alice", "95"], ["Bob", "88"]]
        path = _create_pptx_with_table(tmp_path, "table.pptx", rows)
        fp = FileProcessor(upload_folder=str(tmp_path))
        result = fp.extract_text_from_pptx(path)
        for cell in ["Name", "Score", "Alice", "95", "Bob", "88"]:
            assert cell in result

    def test_empty_presentation(self, tmp_path):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank slide, no shapes
        path = os.path.join(tmp_path, "empty.pptx")
        prs.save(path)
        fp = FileProcessor(upload_folder=str(tmp_path))
        result = fp.extract_text_from_pptx(path)
        assert result == ""

    def test_process_file_routes_pptx(self, tmp_path):
        """process_file() should route .pptx files to extract_text_from_pptx."""
        _create_pptx(tmp_path, "routed.pptx", ["Route test"])
        fp = FileProcessor(upload_folder=str(tmp_path))

        mock_file = MagicMock()
        mock_file.filename = "routed.pptx"
        # File is already in upload_folder (tmp_path), so save is a no-op
        mock_file.save = MagicMock(side_effect=lambda dest: None)

        result = fp.process_file(mock_file)
        assert "Route test" in result

    def test_unsupported_ppt_extension(self, tmp_path):
        """Legacy .ppt files are not supported; process_file returns empty."""
        dummy = os.path.join(tmp_path, "old.ppt")
        with open(dummy, "w") as f:
            f.write("not a real ppt")

        fp = FileProcessor(upload_folder=str(tmp_path))
        mock_file = MagicMock()
        mock_file.filename = "old.ppt"
        # File is already in upload_folder (tmp_path), so save is a no-op
        mock_file.save = MagicMock(side_effect=lambda dest: None)

        result = fp.process_file(mock_file)
        assert result == ""

    def test_multiple_slides(self, tmp_path):
        """Text from multiple slides should all be extracted."""
        prs = Presentation()
        for text in ["Slide 1 content", "Slide 2 content", "Slide 3 content"]:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(5), Inches(1)
            )
            txBox.text_frame.text = text
        path = os.path.join(tmp_path, "multi.pptx")
        prs.save(path)

        fp = FileProcessor(upload_folder=str(tmp_path))
        result = fp.extract_text_from_pptx(path)
        assert "Slide 1 content" in result
        assert "Slide 2 content" in result
        assert "Slide 3 content" in result
